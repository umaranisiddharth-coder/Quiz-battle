from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# Colors
DARK  = RGBColor(0x0a, 0x0a, 0x1a)
PURP  = RGBColor(0x66, 0x7e, 0xea)
PINK  = RGBColor(0xf0, 0x93, 0xfb)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY  = RGBColor(0x99, 0x99, 0xbb)
GREEN = RGBColor(0x43, 0xe9, 0x7b)
GOLD  = RGBColor(0xff, 0xd2, 0x00)
DBLUE = RGBColor(0x12, 0x12, 0x35)
MBLUE = RGBColor(0x18, 0x18, 0x45)

def bg(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = DARK

def box(slide, text, x, y, w, h, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color

def rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()

def bullets(slide, items, x, y, w, h, size=17, color=WHITE, prefix='  '):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = prefix + item
        r.font.size = Pt(size)
        r.font.color.rgb = color

def header(slide, title):
    rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.1), RGBColor(0x10,0x10,0x30))
    box(slide, title, Inches(0.5), Inches(0.1), Inches(12.33), Inches(0.9),
        size=30, color=PURP, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 1 - TITLE
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
rect(sl, Inches(0), Inches(2.2), Inches(13.33), Inches(3.0), RGBColor(0x10,0x10,0x28))

box(sl, 'QUIZ BATTLE GAME', Inches(0.5), Inches(1.0),
    Inches(12.33), Inches(1.2), size=50, color=PURP, bold=True, align=PP_ALIGN.CENTER)

box(sl, 'A fun online quiz game where you can play alone or fight against friends',
    Inches(1.0), Inches(2.3), Inches(11.33), Inches(0.7),
    size=20, color=WHITE, align=PP_ALIGN.CENTER)

box(sl, 'Made with Python, Flask, and HTML/CSS/JavaScript',
    Inches(1.0), Inches(3.1), Inches(11.33), Inches(0.5),
    size=16, color=GRAY, align=PP_ALIGN.CENTER)

# Team names
rect(sl, Inches(0), Inches(5.6), Inches(13.33), Inches(1.9), RGBColor(0x0e,0x0e,0x28))
box(sl, 'Made By', Inches(0.5), Inches(5.65), Inches(12.33), Inches(0.4),
    size=13, color=GRAY, bold=True, align=PP_ALIGN.CENTER)

team = ['Sidhanth Umarane', 'Sanket Yelmate', 'Gopal Solnaki']
for i, name in enumerate(team):
    x = Inches(1.2 + i * 3.7)
    rect(sl, x, Inches(6.1), Inches(3.3), Inches(1.1), MBLUE)
    box(sl, name, x, Inches(6.15), Inches(3.3), Inches(0.6),
        size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    box(sl, 'Developer', x, Inches(6.7), Inches(3.3), Inches(0.4),
        size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2 - TEAM MEMBERS
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
header(sl, 'OUR TEAM')

box(sl, 'Computer Science Project  |  2024-25',
    Inches(0.5), Inches(1.2), Inches(12.33), Inches(0.4),
    size=14, color=GRAY, align=PP_ALIGN.CENTER)

team_data = [
    ('S', 'Sidhanth Umarane',  'Developer',
     'Built the backend server\nCreated the database\nDesigned the game logic'),
    ('S', 'Sanket Yelmate',    'Developer',
     'Built the quiz screens\nCreated battle mode\nDid testing and fixing'),
    ('G', 'Gopal Solnaki',     'Developer',
     'Added all questions\nBuilt the admin panel\nHelped with deployment'),
]
for i, (letter, name, role, work) in enumerate(team_data):
    x = Inches(0.5 + i * 4.2)
    y = Inches(1.8)
    rect(sl, x, y, Inches(3.9), Inches(5.3), DBLUE)
    rect(sl, x, y, Inches(3.9), Inches(0.07), PURP)

    # Big letter avatar
    rect(sl, x + Inches(1.2), y + Inches(0.3), Inches(1.5), Inches(1.5), MBLUE)
    box(sl, letter, x + Inches(1.2), y + Inches(0.35), Inches(1.5), Inches(1.4),
        size=52, color=PURP, bold=True, align=PP_ALIGN.CENTER)

    box(sl, name, x + Inches(0.1), y + Inches(2.0), Inches(3.7), Inches(0.55),
        size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    rect(sl, x + Inches(0.6), y + Inches(2.65), Inches(2.7), Inches(0.45), MBLUE)
    box(sl, role, x + Inches(0.6), y + Inches(2.67), Inches(2.7), Inches(0.4),
        size=13, color=PURP, bold=True, align=PP_ALIGN.CENTER)

    box(sl, work, x + Inches(0.15), y + Inches(3.25), Inches(3.6), Inches(1.5),
        size=13, color=GRAY, align=PP_ALIGN.CENTER)

    box(sl, 'Quiz Battle Game', x + Inches(0.1), y + Inches(4.85), Inches(3.7), Inches(0.4),
        size=12, color=PINK, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 3 - WHAT IS THIS PROJECT
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
header(sl, 'WHAT IS THIS PROJECT?')

box(sl, 'Quiz Battle is a website where you can play quiz games.',
    Inches(0.6), Inches(1.3), Inches(12), Inches(0.6),
    size=20, color=WHITE)

bullets(sl, [
    'You answer 10 questions one by one',
    'Each question has 4 choices - pick the right one',
    'You have 15 seconds to answer each question',
    'You get points for correct answers',
    'Faster answers give you more points',
    'You can play alone, against a computer, or against a friend',
    'Your score is saved and shown on the leaderboard',
    '181 questions on General Knowledge, India, and Computer topics',
], Inches(0.6), Inches(2.0), Inches(12), Inches(5.0), size=19, color=WHITE, prefix='  ->  ')

# ============================================================
# SLIDE 4 - HOW TO PLAY
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
header(sl, 'HOW TO PLAY')

steps = [
    ('Step 1', 'Open the website and type your name'),
    ('Step 2', 'Pick your avatar (emoji picture)'),
    ('Step 3', 'Choose a game mode - Solo, vs Computer, or Battle'),
    ('Step 4', 'Pick a topic - General, India, or Computer'),
    ('Step 5', 'Answer 10 questions before the timer runs out'),
    ('Step 6', 'See your score and check the leaderboard'),
]
for i, (step, desc) in enumerate(steps):
    row = i // 2
    col = i % 2
    x = Inches(0.4 + col * 6.4)
    y = Inches(1.3 + row * 1.8)
    rect(sl, x, y, Inches(6.0), Inches(1.55), DBLUE)
    rect(sl, x, y, Inches(6.0), Inches(0.06), PURP)
    box(sl, step, x + Inches(0.15), y + Inches(0.1), Inches(5.7), Inches(0.45),
        size=15, color=PURP, bold=True)
    box(sl, desc, x + Inches(0.15), y + Inches(0.55), Inches(5.7), Inches(0.85),
        size=16, color=WHITE)

# ============================================================
# SLIDE 5 - GAME MODES
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
header(sl, 'GAME MODES')

modes = [
    ('Solo Quiz', 'Play alone against questions.\nAnswer 10 questions with a timer.\nUse lifelines if you are stuck.'),
    ('VS Computer', 'Play against a computer bot.\nChoose Easy, Medium, or Hard bot.\nSee who scores more points.'),
    ('Battle Mode', 'Play against a real friend.\nShare a room code with your friend.\nBoth answer same questions live.'),
    ('Daily Challenge', 'Special quiz every day.\nSame questions for all players.\nEarn bonus points.'),
]
for i, (title, desc) in enumerate(modes):
    row = i // 2
    col = i % 2
    x = Inches(0.4 + col * 6.4)
    y = Inches(1.3 + row * 2.8)
    rect(sl, x, y, Inches(6.0), Inches(2.5), DBLUE)
    rect(sl, x, y, Inches(6.0), Inches(0.07), PURP)
    box(sl, title, x + Inches(0.15), y + Inches(0.15), Inches(5.7), Inches(0.5),
        size=18, color=PURP, bold=True)
    box(sl, desc, x + Inches(0.15), y + Inches(0.7), Inches(5.7), Inches(1.6),
        size=15, color=WHITE)

# ============================================================
# SLIDE 6 - LIFELINES AND SCORING
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
header(sl, 'LIFELINES AND SCORING')

box(sl, 'Lifelines  (helpers you can use once per game)',
    Inches(0.5), Inches(1.3), Inches(12), Inches(0.5),
    size=18, color=PURP, bold=True)
bullets(sl, [
    '50-50  ->  Removes 2 wrong answers so only 2 choices are left',
    'Skip   ->  Skip a question you do not know',
    'Hint   ->  Shows a small clue about the answer (costs 5 points)',
], Inches(0.5), Inches(1.9), Inches(12), Inches(1.8), size=17, color=WHITE, prefix='  ')

rect(sl, Inches(0.3), Inches(3.8), Inches(12.7), Inches(0.04), RGBColor(0x22,0x22,0x44))

box(sl, 'How Points Work',
    Inches(0.5), Inches(3.95), Inches(12), Inches(0.5),
    size=18, color=PURP, bold=True)
bullets(sl, [
    'Correct answer  ->  10 points',
    'Fast answer bonus  ->  up to 5 extra points',
    'Daily challenge  ->  15 points per correct answer',
    'Hint penalty  ->  minus 5 points',
    '3 correct in a row  ->  Streak bonus shown on screen',
], Inches(0.5), Inches(4.55), Inches(12), Inches(2.5), size=17, color=WHITE, prefix='  ')

# ============================================================
# SLIDE 7 - PLAYER LEVELS
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
header(sl, 'PLAYER LEVELS')

box(sl, 'As you play more games and earn points, your level goes up.',
    Inches(0.5), Inches(1.3), Inches(12), Inches(0.5),
    size=18, color=WHITE)

levels = [
    ('Rookie',  '0 to 99 points',    'You are just starting out'),
    ('Player',  '100 to 299 points', 'You know the basics'),
    ('Pro',     '300 to 599 points', 'You are getting good'),
    ('Expert',  '600 to 999 points', 'You know a lot'),
    ('Legend',  '1000+ points',      'You are the best'),
]
emojis = ['Rookie', 'Player', 'Pro', 'Expert', 'Legend']
lcolors = [GREEN, PURP, GOLD, PINK, WHITE]
for i, (name, pts, desc) in enumerate(levels):
    y = Inches(2.0 + i * 0.95)
    rect(sl, Inches(0.5), y, Inches(12.33), Inches(0.82), DBLUE)
    box(sl, name, Inches(0.7), y + Inches(0.1), Inches(2.5), Inches(0.6),
        size=17, color=lcolors[i], bold=True)
    box(sl, pts, Inches(3.5), y + Inches(0.1), Inches(3.0), Inches(0.6),
        size=16, color=GOLD)
    box(sl, desc, Inches(7.0), y + Inches(0.1), Inches(5.5), Inches(0.6),
        size=15, color=GRAY)

# ============================================================
# SLIDE 8 - ACHIEVEMENTS
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
header(sl, 'ACHIEVEMENTS  (Badges to Collect)')

box(sl, 'You unlock badges when you do special things in the game.',
    Inches(0.5), Inches(1.2), Inches(12), Inches(0.45),
    size=17, color=WHITE)

achs = [
    ('First Blood',   'Play your first game'),
    ('Perfectionist', 'Get 100% correct answers'),
    ('On Fire',       'Answer 5 in a row correctly'),
    ('Veteran',       'Play 10 games'),
    ('High Scorer',   'Reach 500 total points'),
    ('Bot Slayer',    'Beat the computer bot'),
    ('Gladiator',     'Win a battle against a friend'),
    ('Daily Warrior', 'Complete the daily challenge'),
    ('Speed Demon',   'Answer in under 3 seconds'),
    ('Regular',       'Play 5 games'),
]
for i, (name, desc) in enumerate(achs):
    row = i // 5
    col = i % 5
    x = Inches(0.3 + col * 2.55)
    y = Inches(1.8 + row * 2.5)
    rect(sl, x, y, Inches(2.35), Inches(2.2), DBLUE)
    rect(sl, x, y, Inches(2.35), Inches(0.06), PURP)
    box(sl, name, x + Inches(0.1), y + Inches(0.15), Inches(2.15), Inches(0.55),
        size=13, color=PURP, bold=True, align=PP_ALIGN.CENTER)
    box(sl, desc, x + Inches(0.1), y + Inches(0.75), Inches(2.15), Inches(1.2),
        size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9 - TOOLS USED
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
header(sl, 'TOOLS AND TECHNOLOGIES USED')

tools = [
    ('Python',        'The main programming language used to build the game'),
    ('Flask',         'A Python tool that runs the website and handles pages'),
    ('Socket.IO',     'Sends live updates between players during battle mode'),
    ('SQLite',        'A simple database that stores users, questions, and scores'),
    ('HTML and CSS',  'Used to design and style all the web pages'),
    ('JavaScript',    'Makes the quiz interactive - timer, buttons, animations'),
    ('Web Audio API', 'Creates sound effects in the browser without any files'),
    ('Chart.js',      'Shows a graph of your score history on the stats page'),
]
for i, (tool, desc) in enumerate(tools):
    row = i // 2
    col = i % 2
    x = Inches(0.4 + col * 6.4)
    y = Inches(1.3 + row * 1.4)
    rect(sl, x, y, Inches(6.0), Inches(1.2), DBLUE)
    box(sl, tool, x + Inches(0.15), y + Inches(0.08), Inches(5.7), Inches(0.45),
        size=16, color=PURP, bold=True)
    box(sl, desc, x + Inches(0.15), y + Inches(0.55), Inches(5.7), Inches(0.55),
        size=13, color=GRAY)

# ============================================================
# SLIDE 10 - DATABASE
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
header(sl, 'DATABASE  (How Data is Stored)')

box(sl, 'We use SQLite database with 3 tables to store all information.',
    Inches(0.5), Inches(1.2), Inches(12), Inches(0.45),
    size=17, color=WHITE)

tables = [
    ('USERS TABLE', [
        'id  -  unique number for each user',
        'name  -  the username',
        'emoji  -  the avatar they picked',
        'total_score  -  all points earned',
        'games_played  -  how many games played',
    ]),
    ('QUESTIONS TABLE', [
        'id  -  unique number for each question',
        'question  -  the question text',
        'option1 to option4  -  the 4 choices',
        'correct_answer  -  which option is right',
        'difficulty  -  easy, medium, or hard',
        'category  -  general, india, or computer',
        'hint  -  a small clue for the question',
    ]),
    ('GAME RESULTS TABLE', [
        'id  -  unique number for each game',
        'user_id  -  which user played',
        'score  -  points earned in that game',
        'correct_count  -  how many right answers',
        'game_mode  -  solo, battle, or daily',
        'category  -  which topic was played',
        'date  -  when the game was played',
    ]),
]
for i, (tname, cols) in enumerate(tables):
    x = Inches(0.3 + i * 4.3)
    y = Inches(1.8)
    rect(sl, x, y, Inches(4.0), Inches(0.5), PURP)
    box(sl, tname, x + Inches(0.1), y + Inches(0.05), Inches(3.8), Inches(0.4),
        size=14, color=WHITE, bold=True)
    for j, col in enumerate(cols):
        cy = y + Inches(0.5 + j * 0.62)
        c = DBLUE if j % 2 == 0 else RGBColor(0x14,0x14,0x30)
        rect(sl, x, cy, Inches(4.0), Inches(0.58), c)
        box(sl, '  ' + col, x + Inches(0.1), cy + Inches(0.08),
            Inches(3.8), Inches(0.45), size=12, color=WHITE)

# ============================================================
# SLIDE 11 - PAGES IN THE WEBSITE
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
header(sl, 'PAGES IN THE WEBSITE')

pages = [
    ('Home Page',         '/',              'Login with name and pick avatar'),
    ('Quiz Screen',       '/static/quiz',   'Answer questions with timer'),
    ('Battle Mode',       '/battle',        'Fight against friend in real time'),
    ('Leaderboard',       '/leaderboard',   'See top players and your rank'),
    ('My Stats',          '/history',       'See your score chart and history'),
    ('Daily Challenge',   '/daily',         'Special quiz that changes every day'),
    ('Achievements',      '/achievements',  'See all badges you have unlocked'),
    ('Admin Panel',       '/admin',         'Add or remove questions from the game'),
]
for i, (name, route, desc) in enumerate(pages):
    row = i // 2
    col = i % 2
    x = Inches(0.4 + col * 6.4)
    y = Inches(1.3 + row * 1.4)
    rect(sl, x, y, Inches(6.0), Inches(1.2), DBLUE)
    box(sl, name, x + Inches(0.15), y + Inches(0.08), Inches(3.5), Inches(0.45),
        size=16, color=PURP, bold=True)
    box(sl, route, x + Inches(3.8), y + Inches(0.08), Inches(2.0), Inches(0.45),
        size=12, color=GOLD)
    box(sl, desc, x + Inches(0.15), y + Inches(0.55), Inches(5.7), Inches(0.55),
        size=13, color=GRAY)

# ============================================================
# SLIDE 12 - WHAT WE LEARNED
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
header(sl, 'WHAT WE LEARNED')

bullets(sl, [
    'How to build a full website using Python and Flask',
    'How to store and get data from a database (SQLite)',
    'How to make two players play together in real time using Socket.IO',
    'How to design a good looking website with HTML and CSS',
    'How to add sound effects using JavaScript',
    'How to show charts and graphs using Chart.js',
    'How to create an admin panel to manage the game',
    'How to work as a team and divide the work properly',
    'How to test the project and fix bugs',
    'How to prepare a project for deployment on the internet',
], Inches(0.6), Inches(1.3), Inches(12), Inches(5.8), size=18, color=WHITE, prefix='  ->  ')

# ============================================================
# SLIDE 13 - FUTURE IDEAS
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
header(sl, 'FUTURE IDEAS')

bullets(sl, [
    'Tournament Mode  -  4 players compete in rounds to find the winner',
    'Mobile App  -  Make it work on Android phones',
    'More Questions  -  Add Science, History, Sports topics',
    'Voice Quiz  -  Read questions out loud using text to speech',
    'Friend System  -  Add friends and challenge them directly',
    'AI Questions  -  Use AI to create new questions automatically',
    'Online Hosting  -  Put the game on the internet for everyone to play',
    'Email Alerts  -  Send reminder for daily challenge every day',
], Inches(0.6), Inches(1.3), Inches(12), Inches(5.5), size=18, color=WHITE, prefix='  ->  ')

# ============================================================
# SLIDE 14 - THANK YOU
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
rect(sl, Inches(0), Inches(2.0), Inches(13.33), Inches(3.5), RGBColor(0x10,0x10,0x28))

box(sl, 'THANK YOU', Inches(0.5), Inches(1.0), Inches(12.33), Inches(1.2),
    size=52, color=PURP, bold=True, align=PP_ALIGN.CENTER)

box(sl, 'Quiz Battle Game', Inches(0.5), Inches(2.3), Inches(12.33), Inches(0.6),
    size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

box(sl, 'A project by Sidhanth Umarane, Sanket Yelmate, and Gopal Solnaki',
    Inches(0.5), Inches(3.1), Inches(12.33), Inches(0.6),
    size=17, color=GRAY, align=PP_ALIGN.CENTER)

box(sl, 'Built with Python  |  Flask  |  Socket.IO  |  SQLite  |  HTML CSS JS',
    Inches(0.5), Inches(3.8), Inches(12.33), Inches(0.5),
    size=14, color=GRAY, align=PP_ALIGN.CENTER)

box(sl, 'Thank you for watching our presentation!',
    Inches(0.5), Inches(5.5), Inches(12.33), Inches(0.5),
    size=16, color=PINK, align=PP_ALIGN.CENTER)

prs.save('QuizBattle_Simple.pptx')
print('Done! Saved as QuizBattle_Simple.pptx')
print('Total slides:', len(prs.slides))

