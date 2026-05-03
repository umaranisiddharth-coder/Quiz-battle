from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation('QuizBattle_Presentation.pptx')

BG    = RGBColor(0x0a, 0x0a, 0x1a)
PURP  = RGBColor(0x66, 0x7e, 0xea)
PINK  = RGBColor(0xf0, 0x93, 0xfb)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY  = RGBColor(0x88, 0x88, 0xaa)
GOLD  = RGBColor(0xff, 0xd2, 0x00)

def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def txt(slide, text, x, y, w, h, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

# ── UPDATE SLIDE 1: Add team names at bottom ─────────────────
sl = prs.slides[0]
rect(sl, Inches(0), Inches(5.7), Inches(13.33), Inches(1.8), RGBColor(0x10,0x10,0x30))
txt(sl, 'PROJECT TEAM', Inches(0.5), Inches(5.75), Inches(12.33), Inches(0.45),
    size=13, color=PURP, bold=True, align=PP_ALIGN.CENTER)

members = [
    ('Sidhanth Umarane', 'Developer'),
    ('Sanket Yelmate',   'Developer'),
    ('Gopal Solnaki',    'Developer'),
]
for i, (name, role) in enumerate(members):
    x = Inches(1.3 + i * 3.6)
    y = Inches(6.2)
    rect(sl, x, y, Inches(3.2), Inches(1.0), RGBColor(0x18,0x18,0x40))
    tb2 = sl.shapes.add_textbox(x + Inches(0.1), y + Inches(0.08), Inches(3.0), Inches(0.85))
    tf2 = tb2.text_frame
    tf2.word_wrap = False
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r1 = p2.add_run()
    r1.text = name
    r1.font.size = Pt(15)
    r1.font.bold = True
    r1.font.color.rgb = WHITE
    p3 = tf2.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    r2 = p3.add_run()
    r2.text = role
    r2.font.size = Pt(11)
    r2.font.color.rgb = GRAY

# ── ADD TEAM SLIDE ───────────────────────────────────────────
new_sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(new_sl)
rect(new_sl, Inches(0), Inches(0), Inches(13.33), Inches(1.1), RGBColor(0x10,0x10,0x30))
txt(new_sl, 'PROJECT TEAM', Inches(0.5), Inches(0.1), Inches(12.33), Inches(0.9),
    size=32, color=PURP, bold=True, align=PP_ALIGN.CENTER)
rect(new_sl, Inches(4), Inches(1.15), Inches(5.33), Inches(0.04), PURP)
txt(new_sl, 'Computer Science Project  |  Academic Year 2024-25',
    Inches(0.5), Inches(1.25), Inches(12.33), Inches(0.4),
    size=14, color=GRAY, align=PP_ALIGN.CENTER)

members_full = [
    ('Sidhanth Umarane',  'Full Stack Developer', 'Backend  |  Database  |  UI Design'),
    ('Sanket Yelmate',    'Full Stack Developer', 'Frontend  |  Battle Mode  |  Testing'),
    ('Gopal Solnaki',     'Full Stack Developer', 'Questions  |  Admin Panel  |  Deploy'),
]
avatars = ['S', 'S', 'G']
for i, (name, role, contrib) in enumerate(members_full):
    x = Inches(0.5 + i * 4.2)
    y = Inches(1.9)
    rect(new_sl, x, y, Inches(3.9), Inches(4.8), RGBColor(0x12,0x12,0x32))
    rect(new_sl, x, y, Inches(3.9), Inches(0.08), PURP)

    # Avatar circle
    rect(new_sl, x + Inches(1.2), y + Inches(0.25), Inches(1.5), Inches(1.5), RGBColor(0x1a,0x1a,0x50))
    txt(new_sl, avatars[i], x + Inches(1.2), y + Inches(0.3), Inches(1.5), Inches(1.4),
        size=52, color=PURP, bold=True, align=PP_ALIGN.CENTER)

    # Name
    txt(new_sl, name, x + Inches(0.1), y + Inches(2.0), Inches(3.7), Inches(0.6),
        size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Role badge
    rect(new_sl, x + Inches(0.5), y + Inches(2.65), Inches(2.9), Inches(0.45), RGBColor(0x20,0x20,0x55))
    txt(new_sl, role, x + Inches(0.5), y + Inches(2.67), Inches(2.9), Inches(0.4),
        size=13, color=PURP, bold=True, align=PP_ALIGN.CENTER)

    # Contribution
    txt(new_sl, contrib, x + Inches(0.1), y + Inches(3.25), Inches(3.7), Inches(0.5),
        size=12, color=GRAY, align=PP_ALIGN.CENTER)

    rect(new_sl, x + Inches(0.3), y + Inches(3.85), Inches(3.3), Inches(0.03), RGBColor(0x22,0x22,0x55))
    txt(new_sl, 'Quiz Battle Game', x + Inches(0.1), y + Inches(4.0), Inches(3.7), Inches(0.5),
        size=12, color=PINK, align=PP_ALIGN.CENTER)

tagline = 'Built with passion, powered by Python'
txt(new_sl, tagline, Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.4),
    size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ── ADD SCREENSHOTS SLIDE ────────────────────────────────────
sc_sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sc_sl)
rect(sc_sl, Inches(0), Inches(0), Inches(13.33), Inches(1.1), RGBColor(0x10,0x10,0x30))
txt(sc_sl, 'APPLICATION SCREENSHOTS', Inches(0.5), Inches(0.1), Inches(12.33), Inches(0.9),
    size=30, color=PURP, bold=True, align=PP_ALIGN.CENTER)

# Screenshot placeholders with labels (images will be added if files exist)
import os
screenshot_files = [
    ('screenshot1.png', 'Login Screen - Avatar Picker', Inches(0.3), Inches(1.3), Inches(4.0), Inches(2.8)),
    ('screenshot2.png', 'Mode Selection Screen',        Inches(4.6), Inches(1.3), Inches(4.0), Inches(2.8)),
    ('screenshot3.png', 'VS Computer Setup',            Inches(9.0), Inches(1.3), Inches(4.0), Inches(2.8)),
]

labels = [
    ('Login Screen', 'Username input + 32 emoji avatars\nDark gaming theme'),
    ('Mode Selection', 'Solo Quiz / VS Computer / Battle Mode\nPlayer welcome bar with rank'),
    ('VS Computer Setup', 'Category selector + Bot difficulty\nEasy / Medium / Hard bot'),
]

for i, (label, desc) in enumerate(labels):
    x = Inches(0.3 + i * 4.35)
    y = Inches(1.3)
    rect(sc_sl, x, y, Inches(4.1), Inches(3.2), RGBColor(0x12,0x12,0x32))
    rect(sc_sl, x, y, Inches(4.1), Inches(0.06), PURP)

    # Try to add actual screenshot image
    img_file = f'screenshot{i+1}.png'
    if os.path.exists(img_file):
        sc_sl.shapes.add_picture(img_file, x + Inches(0.05), y + Inches(0.1),
                                  Inches(4.0), Inches(2.4))
    else:
        # Placeholder
        rect(sc_sl, x + Inches(0.05), y + Inches(0.1), Inches(4.0), Inches(2.4), RGBColor(0x18,0x18,0x45))
        txt(sc_sl, '[ Screenshot ]', x + Inches(0.05), y + Inches(1.0), Inches(4.0), Inches(0.5),
            size=14, color=RGBColor(0x33,0x33,0x66), align=PP_ALIGN.CENTER)

    txt(sc_sl, label, x + Inches(0.1), y + Inches(2.55), Inches(3.9), Inches(0.4),
        size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(sc_sl, desc, x + Inches(0.1), y + Inches(2.95), Inches(3.9), Inches(0.5),
        size=11, color=GRAY, align=PP_ALIGN.CENTER)

# Second row of screenshots
labels2 = [
    ('Quiz Screen', 'Question + Timer + Lifelines\nStreak display + HUD'),
    ('Battle Mode', 'VS bar with avatars + live scores\nReal-time WebSocket'),
    ('Leaderboard', 'Top players + YOUR rank card\nEmoji avatars + stats'),
]
for i, (label, desc) in enumerate(labels2):
    x = Inches(0.3 + i * 4.35)
    y = Inches(4.7)
    rect(sc_sl, x, y, Inches(4.1), Inches(2.5), RGBColor(0x12,0x12,0x32))
    rect(sc_sl, x, y, Inches(4.1), Inches(0.06), PINK)
    rect(sc_sl, x + Inches(0.05), y + Inches(0.1), Inches(4.0), Inches(1.6), RGBColor(0x18,0x18,0x45))
    txt(sc_sl, '[ Screenshot ]', x + Inches(0.05), y + Inches(0.6), Inches(4.0), Inches(0.5),
        size=13, color=RGBColor(0x33,0x33,0x66), align=PP_ALIGN.CENTER)
    txt(sc_sl, label, x + Inches(0.1), y + Inches(1.75), Inches(3.9), Inches(0.4),
        size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(sc_sl, desc, x + Inches(0.1), y + Inches(2.15), Inches(3.9), Inches(0.4),
        size=11, color=GRAY, align=PP_ALIGN.CENTER)

# ── REORDER: Move team slide to position 2, screenshots to 3 ─
xml_slides = prs.slides._sldIdLst
# Last two added are team and screenshots
team_slide = xml_slides[-2]
ss_slide   = xml_slides[-1]
xml_slides.remove(team_slide)
xml_slides.remove(ss_slide)
xml_slides.insert(1, ss_slide)
xml_slides.insert(1, team_slide)

prs.save('QuizBattle_Final.pptx')
print('Done! Saved as QuizBattle_Final.pptx')
print('Total slides:', len(prs.slides))
for i, sl in enumerate(prs.slides):
    print(f'  Slide {i+1}')
